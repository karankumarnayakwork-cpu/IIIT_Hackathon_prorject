from pptx import Presentation
import os


def create_presentation(title, slides, sources):

    # Ensure output folder exists
    os.makedirs("output", exist_ok=True)

    # Load template
    prs = Presentation("template.pptx")

    # ------------------------------------------------
    # Title Slide
    # ------------------------------------------------

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = "AI Generated Briefing"

    # ------------------------------------------------
    # Content Slides
    # ------------------------------------------------

    for slide_title, points in slides.items():

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_title

        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for point in points:

            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0

    # ------------------------------------------------
    # References Slide
    # ------------------------------------------------

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = "References"

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for src in sources:

        p = text_frame.add_paragraph()
        p.text = src
        p.level = 0

    # ------------------------------------------------
    # Save Presentation
    # ------------------------------------------------

    output_path = "output/briefing.pptx"
    prs.save(output_path)

    return output_path
